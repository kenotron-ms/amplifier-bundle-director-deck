import pytest
from pathlib import Path
from pptx import Presentation
from director_deck.schema import SlideDeck, DesignTokens
from director_deck.pptx_builder import build_pptx, SLIDE_WIDTH_IN, SLIDE_HEIGHT_IN

_EMU_PER_INCH = 914400


@pytest.fixture
def tokens(sample_design_md):
    return DesignTokens.from_design_md(sample_design_md)


@pytest.fixture
def deck(sample_deck_data):
    return SlideDeck.model_validate(sample_deck_data)


@pytest.fixture
def minimal_png(tmp_path) -> Path:
    """A valid 1×1 white PNG."""
    png_bytes = (
        b"\x89PNG\r\n\x1a\n"
        b"\x00\x00\x00\rIHDR"
        b"\x00\x00\x00\x01\x00\x00\x00\x01"
        b"\x08\x02\x00\x00\x00"
        b"\x90wS\xde"
        b"\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N"
        b"\x00\x00\x00\x00IEND\xaeB`\x82"
    )
    p = tmp_path / "tiny.png"
    p.write_bytes(png_bytes)
    return p


class TestBuildPptx:
    def test_slide_count_matches_deck(self, deck, tokens, tmp_path):
        out = tmp_path / "deck.pptx"
        build_pptx(deck, tokens, out)
        prs = Presentation(str(out))
        assert len(prs.slides) == deck.meta.slide_count

    def test_slide_dimensions_are_16_9(self, deck, tokens, tmp_path):
        out = tmp_path / "deck.pptx"
        build_pptx(deck, tokens, out)
        prs = Presentation(str(out))
        width_in = float(prs.slide_width) / _EMU_PER_INCH
        height_in = float(prs.slide_height) / _EMU_PER_INCH
        assert abs(width_in - SLIDE_WIDTH_IN) < 0.02
        assert abs(height_in - SLIDE_HEIGHT_IN) < 0.02

    def test_title_text_appears_in_slide_one(self, deck, tokens, tmp_path):
        out = tmp_path / "deck.pptx"
        build_pptx(deck, tokens, out)
        prs = Presentation(str(out))
        all_text = " ".join(
            shape.text_frame.text
            for shape in prs.slides[0].shapes
            if shape.has_text_frame
        )
        assert "Slide One" in all_text

    def test_enriched_mode_with_real_image(
        self, tokens, tmp_path, minimal_png, sample_deck_data
    ):
        data = {
            **sample_deck_data,
            "meta": {**sample_deck_data["meta"], "slide_count": 1},
            "slides": [
                {
                    **sample_deck_data["slides"][0],
                    "assets": {"image": str(minimal_png), "backdrop": None},
                }
            ],
        }
        enriched_deck = SlideDeck.model_validate(data)
        out = tmp_path / "enriched.pptx"
        build_pptx(enriched_deck, tokens, out, enriched=True)
        prs = Presentation(str(out))
        assert len(prs.slides) == 1

    def test_enriched_mode_missing_asset_does_not_crash(self, deck, tokens, tmp_path):
        """When enriched=True but no asset paths exist, should fall back to placeholder."""
        out = tmp_path / "enriched.pptx"
        build_pptx(deck, tokens, out, enriched=True)
        assert out.exists()
