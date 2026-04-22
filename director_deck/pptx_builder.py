from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from director_deck.schema import DesignTokens, Slide, SlideDeck

# Slide dimensions: 13.33 × 7.5 inches (standard PowerPoint 16:9 widescreen)
SLIDE_WIDTH_IN: float = 13.33
SLIDE_HEIGHT_IN: float = 7.5


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert a #RRGGBB hex string to python-pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_slide(
    prs: Presentation,
    slide: Slide,
    tokens: DesignTokens,
    *,
    enriched: bool = False,
) -> None:
    """Add one populated slide to the Presentation."""
    blank_layout = prs.slide_layouts[6]  # index 6 = blank layout
    sl = prs.slides.add_slide(blank_layout)

    primary_rgb = _hex_to_rgb(tokens.colors.primary)
    accent_rgb = _hex_to_rgb(tokens.colors.accent or "#38BDF8")
    on_surface_rgb = _hex_to_rgb(tokens.colors.on_surface or "#F1F5F9")
    surface_rgb = _hex_to_rgb(tokens.colors.surface or "#1E293B")

    # ── Background fill ──────────────────────────────────────────────────────
    bg_fill = sl.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = primary_rgb

    # ── Title text box — top strip, full width ───────────────────────────────
    title_box = sl.shapes.add_textbox(
        Inches(0.5),
        Inches(0.3),
        Inches(SLIDE_WIDTH_IN - 1.0),
        Inches(1.2),
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.text = slide.title
    para.alignment = PP_ALIGN.LEFT
    run = para.runs[0]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = accent_rgb

    # ── Bullet text box — left 55% of slide, below title ────────────────────
    body_box = sl.shapes.add_textbox(
        Inches(0.5),
        Inches(1.7),
        Inches(SLIDE_WIDTH_IN * 0.55 - 0.5),
        Inches(SLIDE_HEIGHT_IN - 2.3),
    )
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    for i, bullet in enumerate(slide.bullets):
        para = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        para.text = f"\u2192 {bullet}"
        para.space_before = Pt(6)
        run = para.runs[0]
        run.font.size = Pt(16)
        run.font.color.rgb = on_surface_rgb

    # ── Image area — right ~40% of slide ────────────────────────────────────
    img_left = Inches(SLIDE_WIDTH_IN * 0.58)
    img_top = Inches(1.7)
    img_width = Inches(SLIDE_WIDTH_IN * 0.38)
    img_height = Inches(SLIDE_HEIGHT_IN - 2.5)

    if enriched and slide.assets.image:
        img_path = Path(slide.assets.image)
        if img_path.exists():
            sl.shapes.add_picture(
                str(img_path), img_left, img_top, img_width, img_height
            )
            return  # real picture added; skip placeholder

    # Placeholder rectangle (wireframe mode or asset file missing)
    ph = sl.shapes.add_shape(1, img_left, img_top, img_width, img_height)
    ph.fill.solid()
    ph.fill.fore_color.rgb = surface_rgb
    ph.line.color.rgb = accent_rgb
    tf_ph = ph.text_frame
    tf_ph.word_wrap = True
    para_ph = tf_ph.paragraphs[0]
    para_ph.text = slide.image_brief
    para_ph.alignment = PP_ALIGN.CENTER
    run_ph = para_ph.runs[0]
    run_ph.font.size = Pt(11)
    run_ph.font.color.rgb = on_surface_rgb


def build_pptx(
    deck: SlideDeck,
    tokens: DesignTokens,
    output_path: Path,
    *,
    enriched: bool = False,
) -> Path:
    """
    Build a PPTX from a SlideDeck.

    Args:
        deck: Validated SlideDeck.
        tokens: Design tokens parsed from DESIGN.md.
        output_path: Destination for the .pptx file.
        enriched: False = placeholder image areas. True = embed real images where available.

    Returns:
        The path the file was written to.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    for slide in deck.slides:
        _add_slide(prs, slide, tokens, enriched=enriched)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path
